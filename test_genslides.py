import collections 
import collections.abc

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

table_slide_layout = prs.slide_layouts[0]

# Adding a blank slide in out ppt
slide = prs.slides.add_slide(prs.slide_layouts[6])
  
# Adjusting the width !  
x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5) 
  
# Adding tables
shape = slide.shapes.add_table(3, 4, x,  y, cx, cy)

table = shape.table

cell = table.cell(0, 0)

cell.text = 'some'
cell = table.cell(1, 0)

cell.text = 'one'
cell = table.cell(2, 0)

cell.text = 'two'

#charts

slide = prs.slides.add_slide(prs.slide_layouts[5])

chart_data = CategoryChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

title = slide.shapes.title
title.text = "Chart"

#image

img_path = 'three.jpg' 
  
# Creating an Presentation object
  
# Selecting blank slide
blank_slide_layout = prs.slide_layouts[6] 
  
# Attaching slide to ppt
slide = prs.slides.add_slide(blank_slide_layout) 
  
# For margins
left = top = Inches(1) 
  
# adding images
# pic = slide.shapes.add_picture(img_path, left, top)
  
left = Inches(2) 
height = Inches(5) 
  
pic = slide.shapes.add_picture(img_path, left, top, height = height)


prs.save('test.pptx')
