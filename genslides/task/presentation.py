from genslides.task.base import BaseTask
from genslides.task.base import TaskDescription

from genslides.task.text import TextTask
from genslides.task.text import RichTextTask
from genslides.commands.create import CreateCommand


import genslides.utils.reqhelper as ReqHelper
import genslides.utils.request as Requester
import genslides.utils.request as Response
from genslides.utils.chatgptrequester import ChatGPTrequester

import collections
import collections.abc

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt, Cm

import os
from os import listdir
from os.path import isfile, join



class PresentationTask(TextTask):
    def __init__(self, task_info : TaskDescription ) -> None:
        super().__init__(task_info, "Presentation")
        print("Start presentation Task")
        request = self.init + self.prompt + self.endi
        descriptions = self.getResponse(request)
        if len(descriptions) == 0:
            print("Request to ChatGPT")
            chat  = ChatGPTrequester()
            responses = chat.getResponse(request, ['title','description'])
            for response in responses:
                print(response)
                description = {}
                description['prompt'] = "Title of slide is " + response.trgs['title'] +". Slide description is:"+  response.trgs['description'] + "Presentation description is :" + self.prompt
                description['title'] = response.trgs['title']
                descriptions.append(description)
            self.saveRespJson(request, descriptions)
            del chat
        else:
            print("Data from saved file")

        for t in descriptions:
            self.task_creation_result += t + '\n'

        self.prs = Presentation()

        for description in descriptions:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
            
            
            ch_prompt = description
            sents = description.split(".")
            title = sents[0][len("Title of slide is "):]
            print("title=",title)
            # ch_prompt = description['prompt']
            # slide.shapes.title.text = description['title']
            slide.shapes.title.text =title 
            
            self.addChildTask(TaskDescription( prompt= ch_prompt, method= SlideTask, parent= self, helper=self.reqhelper, requester=self.requester, target=slide))
            break

    def completeTask(self):
        if not self.checkChilds():
            return False
        mypath = "output/"
        if not os.path.exists(mypath):
            os.makedirs(mypath)

        name = mypath + "text.pptx"
        print("Save prsenetation=",name)
        self.prs.save(name)
        self.is_solved = True

        return True     

    # def __init__(self,  parent, reqhelper : ReqHelper, requester :Requester, description) -> None:
    #     super().__init__(reqhelper, requester, type='Presentation', prompt=description, parent=parent)
    #     self.prs = Presentation()
    #     self.description = description
    #     title_slide_layout = self.prs.slide_layouts[0]
    #     slide = self.prs.slides.add_slide(title_slide_layout)
    #     title = slide.shapes.title
    #     subtitle = slide.placeholders[1]
    #     print("Init presentation")
    #     presentation_prompt = self.reqhelper.getPrompt('Presentation', self.description)
    #     index = 0
    #     while len(self.responselist) == 0 and index < 3:
    #         print("Request to chatgpt about slides")
    #         if index > 0:
    #             info_prompt = self.reqhelper.getPrompt('Information', self.description)

    #             info_response = self.requester.getResponse(info_prompt)
    #             search_result = ""
    #             for pack in info_response:
    #                 if(pack.type == 'Search'):
    #                     print('Doing search')
    #             comb_prompt = self.reqhelper.getPrompt('Combinator', search_result)
    #             info_response = self.requester.getResponse(comb_prompt)
    #             for pack in info_response:
    #                 if(pack.type == 'Text'):
    #                     print('Add data to description')

                

    #         response = self.requester.getResponse(presentation_prompt)
    #         for pack in response:
    #             if pack.type == 'Slide':
    #                 self.responselist.append(pack)


class SlideTask(TextTask):
    def __init__(self, task_info : TaskDescription) -> None:
        super().__init__(task_info, "Slide")
    # def __init__(self, reqhelper: ReqHelper, requester: Requester, prompt=None, parent=None, method=None) -> None:
        # super().__init__(reqhelper, requester, "Slide", prompt, parent, method)
        print("Create slide with prompt:\n" + self.prompt)
        request = self.init + self.prompt + self.endi
        descriptions = self.getResponse(request)
        slide_part_prompts =[]
        if len(descriptions) == 0:
            self.task_creation_result += "Request to ChatGPT:\n"
            # print(request)
            chat  = ChatGPTrequester()
            responses = chat.getResponse(request, ['type','description','position'])
            index = 0
            for response in responses:
                description = {}
                description['id'] = index
                index += 1
                description['solved'] = False
                description['type'] = response.trgs['type']
                description['description'] = response.trgs['description']
                description['position'] = response.trgs['position']
                # print(description)
                init_sent = "unknown"
                if description['type'] == 'text':
                    init_sent = "detailed and enriched text for slide part\"" + description['description'] + "\". "
                elif description['type'] == 'table':
                    init_sent = "list of table cells for slide part\"" + description['description'] + "\" formatted in json as: cell row, cell coll, cell text. "
                elif description['type'] == 'plot':
                    init_sent = "list of data plots for slide part\"" + description['description'] + "\" formatted in json as: data name, array of data series. "
                elif description['type'] == 'chart':
                    init_sent = "list of data charts for slide part\"" + description['description'] + "\" formatted in json as: data name, data value. "
                elif description['type'] == 'image':
                    init_sent = "image description for slide part\"" + description['description'] + "\". "
                sent = init_sent + self.prompt
                description['prompt'] = sent
                slide_part_prompts.append(sent)
                descriptions.append(description)
            self.saveRespJson(request, descriptions)
            del chat
        else:
            self.task_creation_result += "Data from saved file:\n"
        
        for t in descriptions:
            self.task_creation_result += str(t) + '\n'

        left = Cm(2) 
        top = Cm(3)
        width = Cm(10)
        height = Cm(18)

        # self.paragraph.text = "test"
        # self.text_paragraph = self.paragraph.text
        # print("text_paragraph=",self.text_paragraph)

        index = 0
        # self.paragraph_list = []
        for description in descriptions:


            if index % 2 == 1:
                left = Cm(14)
            else:
                left = Cm(2)
                top = Cm(3) + index * Cm(5)

            txBox = self.target.shapes.add_textbox(left, top, width, height)
            pointer = self.target.shapes.index(txBox)
            tf = txBox.text_frame
            tf.word_wrap = True
            self.paragraph = tf.paragraphs[0]
            self.paragraph.font.size = Pt(8)
            print("target shapes=", self.target.shapes[index])

            self.addChildTask(TaskDescription(
                prompt=description['prompt'], method=RichTextTask, parent=self, helper=self.reqhelper, requester=self.requester, target=self.insertText, id=pointer))
            index += 1
            # if index > 2:
            #     break

    def insertText(self, text: str, id):
        # self.paragraph.text = text
        self.target.shapes[id].text_frame.paragraphs[0].text = text
        # self.parent = parentext_paragrapht
        # self.description = description
        # print("Init slide")
    def completeTask(self):
        if not self.checkChilds():
            return False
        self.is_solved = True       

        return True     

